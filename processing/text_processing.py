import PyPDF2
import re
import os
import nltk
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from Tokenization import preprocess_text
from nltk.tokenize import word_tokenize


# Diretório temporário para download
TEMP_DOWNLOAD_FOLDER = "temp_download"
os.makedirs(TEMP_DOWNLOAD_FOLDER, exist_ok=True)

# Verificando se o tokenizador punkt está baixado
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Verificando se o recurso punkt_tab está baixado
try:
    nltk.data.find('tokenizers/punkt_tab')
except LookupError:
    nltk.download('punkt_tab')


def extract_text_from_pdf(file_path):
    """Extrai texto de um arquivo PDF."""
    text = ""
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
    except Exception as e:
        print(f"Erro ao extrair texto em PDF: {e}")

    return text


def extract_text_from_docx(file_path):
    """Extrai texto de um arquivo DOCX."""
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Erro ao extrair texto do DOCX: {e}")

    return text


def extract_text_from_txt(file_path):
    """Extrai texto de um arquivo TXT."""
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()
    except Exception as e:
        print(f"Erro ao extrair texto do TXT: {e}")

    return text


def extract_text_from_xlsx(file_path):
    """Extrai texto de um arquivo XLSX."""
    text = ""
    try:
        workbook = load_workbook(filename=file_path, read_only=True)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.rows:
                row_text = " ".join(str(cell.value) for cell in row if cell.value is not None)
                text += row_text + "\n"
    except Exception as e:
        print(f"Erro ao extrair texto do XLSX: {e}")

    return text


def extract_text_from_ppt(file_path):
    """Extrai texto de um arquivo PPT."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
                    text += "\n"
    except Exception as e:
        print(f"Erro ao processar o arquivo PPT: {e}")

    return text


def extract_text(file_path):
    """Função genérica para extrair texto com base na extensão do arquivo."""
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".txt"):
        return extract_text_from_txt(file_path)
    elif file_path.endswith(".xlsx"):
        return extract_text_from_xlsx(file_path)
    elif file_path.endswith(".pptx") or file_path.endswith(".ppt"):
        return extract_text_from_ppt(file_path)
    else:
        print(f"Formato de arquivo não suportado para: {file_path}")
        return ""


def process_and_tokenize_file(file_path):
    """Extrai texto de um arquivo e o tokeniza."""
    text = extract_text(file_path)

    if text:
        tokens = preprocess_text(text) # Retorna tokens
        # Mostra os 20 primeiros tokens.
        print(f"Texto tokenizado de '{os.path.basename(file_path)}':\n{tokens[:20]}...\n")
        return os.path.basename(file_path), tokens
    else:
        print(f"Não foi possível extrair texto de '{os.path.basename(file_path)}'.")
        return os.path.basename(file_path), None


def cleanup_temp_folder():
    """Limpa o diretório temporário de download."""
    for filename in os.listdir(TEMP_DOWNLOAD_FOLDER):
        file_path = os.path.join(TEMP_DOWNLOAD_FOLDER, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
        except Exception as e:
            print(f"Erro ao remover {file_path}: {e}")
    os.rmdir(TEMP_DOWNLOAD_FOLDER)
    print(f"Diretório temporário '{TEMP_DOWNLOAD_FOLDER}' limpo.")


def remove_special_characters(text):
    """Remove caracteres especiais do texto, preservando letras (incluindo acentuadas),
    números, espaços e sinais diacríticos comuns em inglês."""
    text = re.sub(r"[^a-zA-Z0-9áàâãéèêíïóôõúüçñÁÀÂÃÉÈÊÍÏÓÔÕÚÜÇÑ\s\-\']", "", text)
    return text


def convert_to_lowercase(text):
    """Converte o texto para minúsculas."""
    return text.lower()


def tokenize_text(text):
    """Tokeniza o texto em palavras."""
    return word_tokenize(text)


def preprocess_text(text):
    """Realiza o pré-processamento completo do texto."""
    tokens = tokenize_text(convert_to_lowercase(remove_special_characters(text)))
    return tokens